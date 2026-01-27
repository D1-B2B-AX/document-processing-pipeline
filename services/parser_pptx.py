from pptx import Presentation
import io

# 이 함수 하나로 모든 라우터가 PPTX를 처리합니다.
def parse_pptx_content(file_content: bytes) -> str:
    prs = Presentation(io.BytesIO(file_content))
    markdown_output = []
    
    for i, slide in enumerate(prs.slides):
        markdown_output.append(f"## Slide {i+1}")
        
        for shape in slide.shapes:
            # 1. 텍스트 상자
            if hasattr(shape, "text") and shape.text.strip():
                markdown_output.append(shape.text.strip())
            
            # 2. 표(Table) 처리
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text_frame.text.strip() for cell in row.cells]
                    markdown_output.append("| " + " | ".join(row_text) + " |")
        
        markdown_output.append("\n---\n")
            
    return "\n".join(markdown_output)