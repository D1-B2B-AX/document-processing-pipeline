from fastapi import APIRouter, UploadFile, File, HTTPException
from pptx import Presentation
from openai import OpenAI
import io
import re
import os

client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))
router = APIRouter()

# =========================================================
# [키워드 설정] (사용자님의 코드를 그대로 가져왔습니다)
# =========================================================
EXCLUDE_KEYWORDS = [
    "유사", "사례", "실적", "reference", "case", "history", "result",
    "강사프로필", "수행실적", "제안사", "회사소개"
]

OVERVIEW_KEYWORDS = [
    "과정 소개", "과정소개", "과정 개요", "과정개요", 
    "교육 소개", "교육소개", "교육 개요", "교육개요", 
    "개요", "소개", "overview", "summary", "요약", "제안 배경", "기획 의도",
    "목표", "대상" 
]

CURRICULUM_KEYWORDS = [
    "커리큘럼", "세부과정", "교육과정", "교육내용", "모듈구성", 
    "상세과정", "프로그램", "module", "schedule", "curriculum",
    "모듈", "구성", "일정", "방법", "contents", "agenda", "syllabus",
    "1일차", "2일차", "1h", "2h", "time" 
]

# =========================================================
# [기능 1] 텍스트 추출 (재귀 + Duck Typing)
# 텍스트가 그룹(Group) 안에 있어도 무조건 꺼내는 핵심 로직입니다.
# =========================================================
def normalize(text):
    return re.sub(r'\s+', '', str(text).lower())

def get_text_from_shape_recursive(shape):
    """도형, 표, 그룹 내부를 가리지 않고 텍스트를 추출합니다."""
    text_parts = []
    try:
        # 1. 텍스트 박스
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            text_parts.append(shape.text.strip())
        
        # 2. 표 (Table)
        if hasattr(shape, "table") and shape.table:
            for row in shape.table.rows:
                row_cells = [c.text.replace('\n', ' ').strip() for c in row.cells if c.text.strip()]
                if row_cells:
                    text_parts.append(f"| {' | '.join(row_cells)} |")
        
        # 3. 그룹 (재귀 탐색) - 사용자님이 원하신 재귀 로직 적용
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                text_parts.extend(get_text_from_shape_recursive(child))
    except:
        pass
    return text_parts

def extract_text_from_slide(slide):
    """슬라이드 전체 텍스트 추출"""
    all_texts = []
    
    # 제목 처리 (Visual Title 로직 대신, 안전하게 객체 속성 확인)
    try:
        if slide.shapes.title and slide.shapes.title.text.strip():
            all_texts.append(f"### {slide.shapes.title.text.strip()}")
    except:
        pass 

    # 본문 처리 (제목 제외)
    for shape in slide.shapes:
        try:
            if slide.shapes.title and shape == slide.shapes.title:
                continue
        except:
            pass
        all_texts.extend(get_text_from_shape_recursive(shape))
        
    return "\n".join(all_texts)

# =========================================================
# [기능 2] 슬라이드 분류 (사용자 로직 + 내용 기반 보완)
# =========================================================
def classify_slide_by_content(full_text):
    """
    제목 위치(Top)에 의존하지 않고, 텍스트 내용을 보고 분류합니다.
    (그룹 안에 제목이 숨어있을 때도 작동하기 위함)
    """
    norm_text = normalize(full_text)
    
    for key in EXCLUDE_KEYWORDS:
        if normalize(key) in norm_text: return "EXCLUDE"
    for key in OVERVIEW_KEYWORDS:
        if normalize(key) in norm_text: return "OVERVIEW"
    for key in CURRICULUM_KEYWORDS:
        if normalize(key) in norm_text: return "CURRICULUM"
    
    return "OTHER"

# =========================================================
# [기능 3] LLM 변환 (표 강제 + 시간 보존)
# =========================================================
def generate_rag_markdown(filename, course_idx, overview_text, curriculum_text):
    if len(curriculum_text) < 30: return None

    # 토큰 제한 안전장치
    safe_curriculum = curriculum_text[:20000]

    prompt = f"""
    당신은 '기업 교육 제안서 분석 전문가'입니다.
    제공된 Raw Text를 분석하여 RAG 검색에 최적화된 **Clean Markdown** 포맷으로 변환하십시오.

    [Input Source]
    - File: {filename}
    - Context (개요): {overview_text[:3000]}
    - Content (커리큘럼): {safe_curriculum}

    [Output Rules - Strict]
    1. **Metadata**: 문서 최상단에 `> **Keywords**: ...` 형식으로 핵심 키워드(대상, 주제, 툴 등) 나열.
    2. **Table Formatting (필수)**: 
       - 커리큘럼의 상세 일정, 모듈 구성은 **반드시 Markdown Table**로 작성할 것.
       - 예시: | 모듈명 | 시간 | 주요 내용 | 교육 방법 |
    3. **Time Preservation**: '1H', '2시간', '09:00~18:00' 등 시간 정보는 **절대 삭제 금지**.
    4. **Filtering**: 강사 프로필, 회사 홍보 등 커리큘럼과 무관한 내용은 삭제.
    5. **No Chit-chat**: 서론 없이 결과 Markdown만 출력.
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            temperature=0
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"LLM Error: {e}")
        return None

# =========================================================
# [Endpoint] 메인 핸들러
# =========================================================
@router.post("/parse")
async def parse_curriculum(file: UploadFile = File(...)):
    print(f"🚀 Processing: {file.filename}")
    
    content = await file.read()
    try:
        prs = Presentation(io.BytesIO(content))
    except Exception as e:
        raise HTTPException(status_code=400, detail="Invalid PPTX file")

    courses = [] 
    current_course = {'overview': [], 'curriculum': []}
    
    # 1. 슬라이드 순회 (사용자님의 로직 반영: OVERVIEW마다 과정 분리)
    for i, slide in enumerate(prs.slides):
        # 재귀함수로 텍스트 추출 (그룹 내부 포함)
        full_text = extract_text_from_slide(slide)
        
        # 내용 기반 분류
        slide_type = classify_slide_by_content(full_text)
        
        if slide_type == "EXCLUDE": 
            continue

        if slide_type == "OVERVIEW":
            # [중요] 새로운 개요가 나오면 이전 과정을 저장하고 리셋 (사용자 로직)
            if current_course['curriculum']: 
                courses.append(current_course)
                current_course = {'overview': [], 'curriculum': []}
            current_course['overview'].append(full_text)

        elif slide_type == "CURRICULUM":
            current_course['curriculum'].append(full_text)
            
        # 분류가 안 된 슬라이드(OTHER)라도 텍스트가 길면 커리큘럼으로 간주 (안전장치)
        elif slide_type == "OTHER" and len(full_text) > 50:
             current_course['curriculum'].append(full_text)

    # 마지막에 남은 과정 추가
    if current_course['curriculum']:
        courses.append(current_course)

    print(f"📊 감지된 과정(Courses) 수: {len(courses)}개")

    # 2. LLM 변환 및 결과 생성
    results = []
    for idx, course in enumerate(courses):
        full_overview = "\n\n".join(course['overview'])
        full_curriculum = "\n\n".join(course['curriculum'])
        
        md_content = generate_rag_markdown(file.filename, idx+1, full_overview, full_curriculum)
        
        if md_content:
            # [수정됨] 파일명 새니타이징(특수문자 제거) 로직 삭제 -> 원본 파일명 유지
            base_name = os.path.splitext(file.filename)[0]
            
            # 과정이 여러 개일 때만 뒤에 번호 붙임, 하나면 깔끔하게 원본명 사용
            suffix = f"_Course_{idx+1}" if len(courses) > 1 else "_Parsed"
            suggested_filename = f"{base_name}{suffix}.md"
            
            results.append({
                "course_index": idx + 1,
                "suggested_filename": suggested_filename,
                "markdown": md_content
            })

    return {
        "domain": "curriculum",
        "original_filename": file.filename,
        "parsed_courses": results,
        "count": len(results)
    }